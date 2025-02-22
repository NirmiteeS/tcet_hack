import customtkinter as ctk
from tkinter import filedialog, messagebox
from gtts import gTTS  # For text-to-speech (Google TTS)
import os
import pygame  # For audio playback
from PyPDF2 import PdfReader  # For PDF text extraction
from pptx import Presentation  # For PPT text extraction
from docx import Document  # For DOCX text extraction
import threading  # For running tasks in a separate thread
import speech_recognition as sr  # For speech-to-text
import google.generativeai as genai  # Gemini API
from fpdf import FPDF  # For generating PDFs

# Initialize pygame for audio playback
pygame.mixer.init()

# Initialize Gemini API
genai.configure(api_key="AIzaSyC5kiR_7pP0x1ft-Fd2mqgakXwl_D7-Kl0")  # Replace with your Gemini API key
model = genai.GenerativeModel('gemini-pro')  # Use the Gemini Pro model

# Global variables for speech-to-text
is_recording = False

# Function to upload a PPT, PDF, or DOCX and convert it to audio (already working)
def upload_file():
    # Open file dialog to select PPT, PDF, or DOCX
    file_path = filedialog.askopenfilename(filetypes=[("PPT Files", "*.pptx"), ("PDF Files", "*.pdf"), ("DOCX Files", "*.docx")])
    if file_path:
        # Convert the file to audio
        audio_filename = convert_to_audio(file_path)
        if audio_filename:
            # After conversion, show the audio file and options to play and download
            show_audio_popup(audio_filename)
        else:
            messagebox.showerror("Error", "Failed to convert file to audio.")

# Function to extract text from PDF
def extract_text_from_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text
    except Exception as e:
        messagebox.showerror("Error", f"Could not extract text from PDF. Error: {e}")
        return None

# Function to extract text from PPT
def extract_text_from_ppt(file_path):
    try:
        presentation = Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        messagebox.showerror("Error", f"Could not extract text from PPT. Error: {e}")
        return None

# Function to extract text from DOCX
def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        messagebox.showerror("Error", f"Could not extract text from DOCX. Error: {e}")
        return None

# Function to convert the extracted text into audio
def convert_to_audio(file_path):
    text = ""
    
    # Extract text based on file type
    if file_path.endswith(".pdf"):
        text = extract_text_from_pdf(file_path)
    elif file_path.endswith(".pptx"):
        text = extract_text_from_ppt(file_path)
    elif file_path.endswith(".docx"):
        text = extract_text_from_docx(file_path)
    
    if not text:
        return None  # If no text was extracted, return None
    
    # Convert the extracted text to speech using gTTS
    audio_filename = "converted_audio.mp3"
    tts = gTTS(text=text, lang='en')
    tts.save(audio_filename)

    return audio_filename

# Function to play audio using pygame
def play_audio(audio_filename):
    try:
        if os.path.exists(audio_filename):
            pygame.mixer.music.load(audio_filename)
            pygame.mixer.music.play()
        else:
            messagebox.showerror("Error", "Audio file not found.")
    except Exception as e:
        messagebox.showerror("Error", f"Could not play the audio. Error: {e}")

# Function to stop audio playback and unload the file
def stop_audio():
    pygame.mixer.music.stop()
    pygame.mixer.music.unload()  # Unload the audio file to release it

# Function to show the audio popup window
def show_audio_popup(audio_filename):
    # Create a popup window for audio options
    audio_window = ctk.CTkToplevel()
    audio_window.title("Audio File Options")
    audio_window.geometry("400x200")

    # Play button for the audio file
    play_btn = ctk.CTkButton(audio_window, text="Play Audio", command=lambda: play_audio(audio_filename), corner_radius=8)
    play_btn.pack(pady=10)

    # Stop button for the audio file
    stop_btn = ctk.CTkButton(audio_window, text="Stop Audio", command=stop_audio, corner_radius=8)
    stop_btn.pack(pady=10)

    # Download button
    def download_audio():
        # Stop audio playback and unload the file before moving it
        stop_audio()
        
        # Ask for download location and save audio file
        download_location = filedialog.asksaveasfilename(defaultextension=".mp3", filetypes=[("MP3 Files", "*.mp3")])
        if download_location:
            try:
                # Move the audio file to the selected location
                os.replace(audio_filename, download_location)
                messagebox.showinfo("Download", f"Audio saved at: {download_location}")
            except Exception as e:
                messagebox.showerror("Error", f"Could not save the audio. Error: {e}")

    download_btn = ctk.CTkButton(audio_window, text="Download Audio", command=download_audio, corner_radius=8)
    download_btn.pack(pady=10)

# Function to convert audio to text using Gemini API
def convert_audio_to_text(file_path):
    try:
        # Use SpeechRecognition to extract text from audio
        recognizer = sr.Recognizer()
        with sr.AudioFile(file_path) as source:
            audio = recognizer.record(source)
            text = recognizer.recognize_google(audio)
            return text
    except Exception as e:
        messagebox.showerror("Error", f"Could not convert audio to text. Error: {e}")
        return None

# Function to generate PDF from text using FPDF
def generate_pdf(text, output_filename):
    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, text)
        pdf.output(output_filename)
        messagebox.showinfo("Success", f"PDF saved as {output_filename}")
    except Exception as e:
        messagebox.showerror("Error", f"Could not generate PDF. Error: {e}")

# Function to generate text file from text
def generate_text_file(text, output_filename):
    try:
        with open(output_filename, "w") as f:
            f.write(text)
        messagebox.showinfo("Success", f"Text file saved as {output_filename}")
    except Exception as e:
        messagebox.showerror("Error", f"Could not generate text file. Error: {e}")

# Function to upload audio file and convert it to PDF or text file
def upload_audio_and_convert():
    # Open file dialog to select an audio file
    file_path = filedialog.askopenfilename(filetypes=[("Audio Files", "*.wav *.mp3")])
    if file_path:
        # Convert audio to text
        text = convert_audio_to_text(file_path)
        if text:
            # Ask the user to choose between PDF and text file
            file_type = messagebox.askquestion("File Type", "Do you want to save as PDF?")
            if file_type == "yes":
                output_filename = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")])
                if output_filename:
                    generate_pdf(text, output_filename)
            else:
                output_filename = filedialog.asksaveasfilename(defaultextension=".txt", filetypes=[("Text Files", "*.txt")])
                if output_filename:
                    generate_text_file(text, output_filename)

# Function to start recording speech and convert it to text using Gemini API
def start_recording():
    global is_recording
    is_recording = True
    threading.Thread(target=record_audio).start()

# Function to stop recording
def stop_recording():
    global is_recording
    is_recording = False

# Function to record audio and convert it to text using Gemini API
def record_audio():
    global is_recording
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        recognizer.adjust_for_ambient_noise(source)
        while is_recording:
            try:
                print("Listening...")
                audio = recognizer.listen(source, phrase_time_limit=5)  # Listen for 5 seconds at a time
                text = recognizer.recognize_google(audio)
                text_entry.insert("end", text + " ")  # Append recognized text to the text box
            except sr.UnknownValueError:
                print("Could not understand audio")
            except sr.RequestError as e:
                print(f"Could not request results; {e}")

# Create the main audio-text window
def create_audio_text_window():
    # Create the window for Audio to Text features
    window = ctk.CTk()
    window.title("Audio to Text Converter")
    window.geometry("600x400")

    # Upload button for PPT/PDF/DOCX to audio
    upload_btn = ctk.CTkButton(window, text="Upload PPT/PDF/DOCX", command=upload_file, corner_radius=8)
    upload_btn.pack(pady=20)

    # Upload audio button
    upload_audio_btn = ctk.CTkButton(window, text="Upload Audio", command=upload_audio_and_convert, corner_radius=8)
    upload_audio_btn.pack(pady=10)

    # Start Recording button
    start_recording_btn = ctk.CTkButton(window, text="Start Recording", command=start_recording, corner_radius=8)
    start_recording_btn.pack(pady=10)

    # Stop Recording button
    stop_recording_btn = ctk.CTkButton(window, text="Stop Recording", command=stop_recording, corner_radius=8)
    stop_recording_btn.pack(pady=10)

    # Text entry for displaying converted text
    global text_entry
    text_entry = ctk.CTkTextbox(window, width=500, height=100, corner_radius=8)
    text_entry.pack(pady=20)

    # Close button to close the window
    close_btn = ctk.CTkButton(window, text="Close", command=window.destroy, corner_radius=8)
    close_btn.pack(pady=20)

    window.mainloop()

# Run the main window
if __name__ == "__main__":
    create_audio_text_window()